#!/usr/bin/env python3
"""Convert existing SVG layouts to real editable PowerPoint objects.

中文：文字、形状和箭头转成 PPT 原生对象，独立嵌入原帧和透明渲染图。
人体网格渲染仍是位图，不声称人体表面变成可编辑矢量网格。
另提供不含实验指标的纯概念示意模板。
"""
from pathlib import Path
import base64
import io
import json
import sys
import xml.etree.ElementTree as ET
import zipfile

ROOT=Path(__file__).resolve().parent
EDIT=ROOT/'editable'
sys.path.insert(0,str(EDIT/'_python_deps'))
from PIL import Image, ImageFont
from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.dml import MSO_LINE_DASH_STYLE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.util import Inches, Pt
import compose_teaser as layout

FONT='/usr/share/fonts/truetype/liberation2/LiberationSans-Regular.ttf'
FONT_BOLD='/usr/share/fonts/truetype/liberation2/LiberationSans-Bold.ttf'
EXTRA_LABELS={
    'Streaming reconstruction across shot changes':'跨镜头切换的流式重建',
    'Conceptual illustration — not an experimental result':'概念示意图，并非实验结果',
    'ONE SHARED WORLD':'统一世界坐标系',
    'Generated scene / illustration placeholder':'生成式场景或插画占位区',
    'Replace this panel with your conceptual artwork':'将此区域替换为你的概念绘图',
    'Illustrated monocular input stream':'示意性的单目输入时间流',
    'Input illustration placeholder':'输入示意图占位区',
    'Keep the same person colors across shots':'人物颜色在各镜头间保持一致',
    'Time':'时间',
}
layout.TRANSLATIONS.update(EXTRA_LABELS)


def conceptual_template():
    # Pure vector template: deliberately no fake reconstruction or metrics.
    layout.OUT=EDIT
    s=layout.SVG(1800,880,'Shot3R conceptual layout / 生成式概念图的可编辑排版模板')
    s.parts[2]='<desc>纯概念排版模板，并非实验输出。可替换为明确标注的生成式插画。三镜头仅用于说明任务，不代表已有三镜头结果。</desc>'
    s.text(30,52,'Shot3R',43,layout.INK,700)
    s.text(250,49,'Streaming reconstruction across shot changes',30,layout.INK,500)
    s.line(30,80,1770,80)
    s.text(1770,114,'Conceptual illustration — not an experimental result',22,layout.MUTED,anchor='end')
    s.rect(30,139,1740,397,'#f5f8fa','#d5dfe5',1,rx=12)
    s.text(68,181,'ONE SHARED WORLD',23,layout.INK,600)
    for i,col in enumerate([layout.CORAL,layout.TEAL,layout.PURPLE]):
        s.dot(1330+i*132,177,col,7)
        s.text(1348+i*132,185,f'Person {i+1}',20,layout.INK)
    s.text(900,305,'Generated scene / illustration placeholder',34,layout.MUTED,anchor='middle')
    s.text(900,349,'Replace this panel with your conceptual artwork',24,layout.MUTED,anchor='middle')
    s.text(900,399,'Keep the same person colors across shots',23,layout.TEAL,600,anchor='middle')
    # Camera icons are schematic, not estimated camera poses.
    for x,col in [(188,layout.BLUE),(889,layout.AMBER),(1577,layout.PURPLE)]:
        s.rect(x-24,450,48,31,'white',col,2,rx=4)
        for xa,ya,xb,yb in [(x+24,457,x+46,448),(x+46,448,x+46,483),(x+46,483,x+24,475)]:
            s.line(xa,ya,xb,yb,col,2)
        s.line(x,488,x,531,col,1.5,'5 5')
    s.text(30,587,'Illustrated monocular input stream',25,layout.INK,600)
    for i,col in enumerate([layout.BLUE,layout.AMBER,layout.PURPLE]):
        x=30+i*585
        s.rect(x,612,570,202,'#f5f8fa','#d5dfe5',1,rx=9)
        s.rect(x+1,613,568,4,col)
        s.text(x+285,652,f'Shot {i+1}',24,col,600,anchor='middle')
        s.text(x+285,729,'Input illustration placeholder',24,layout.MUTED,anchor='middle')
    for x in [607,1192]:
        s.rect(x-48,765,96,30,'white')
        s.text(x,787,'Shot cut',20,layout.AMBER,600,anchor='middle')
    s.line(30,844,1767,844,layout.MUTED,1.5,arrow=True)
    s.text(1770,873,'Time',20,layout.MUTED,anchor='end')
    s.save('Shot3R_concept_template')
    return EDIT/'Shot3R_concept_template.svg'


def rgb(value):
    return RGBColor.from_string(value.lstrip('#'))


def build(source,destination,concept=False):
    parser=ET.XMLParser(target=ET.TreeBuilder(insert_comments=True))
    tree=ET.parse(source,parser=parser)
    root=tree.getroot()
    _,_,vw,vh=map(float,root.attrib['viewBox'].split())
    prs=Presentation()
    prs.slide_width=Inches(13.5)
    prs.slide_height=Inches(13.5*vh/vw)
    unit=prs.slide_width/vw
    px=lambda v: round(float(v)*unit)
    pt=lambda v: Pt(float(v)*13.5*72/vw)
    slide=prs.slides.add_slide(prs.slide_layouts[6])
    counts={'text_boxes':0,'pictures':0,'vector_shapes':0,'connectors':0}
    labels=[]
    last_asset=''

    def line_style(shape,stroke,width=1,dash=None):
        if stroke in ['none',None]:
            shape.line.fill.background()
        else:
            shape.line.color.rgb=rgb(stroke)
            shape.line.width=pt(width)
            if dash:shape.line.dash_style=MSO_LINE_DASH_STYLE.DASH

    def fill_style(shape,fill):
        if fill in ['none',None]:shape.fill.background()
        else:
            shape.fill.solid()
            shape.fill.fore_color.rgb=rgb(fill) if fill.startswith('#') else RGBColor(255,255,255)

    for e in root:
        if e.tag is ET.Comment:
            if (e.text or '').strip().startswith('Asset:'):
                last_asset=e.text.strip().split('Asset:',1)[1].strip()
            continue
        kind=e.tag.rsplit('}',1)[-1]
        a=e.attrib
        if kind in ['title','desc','defs']:continue
        if kind=='rect':
            w=vw if a['width']=='100%' else float(a['width'])
            h=vh if a['height']=='100%' else float(a['height'])
            x,y=float(a.get('x',0)),float(a.get('y',0))
            shape=slide.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE if float(a.get('rx',0)) else MSO_SHAPE.RECTANGLE,
                    px(x),px(y),px(w),px(h))
            if shape.auto_shape_type==MSO_SHAPE.ROUNDED_RECTANGLE:
                shape.adjustments[0]=min(.5,float(a['rx'])/max(1,min(w,h)))
            fill_style(shape,a.get('fill','none'))
            line_style(shape,a.get('stroke','none'),a.get('stroke-width',1))
            shape.name=f'Panel_{counts["vector_shapes"]+1}'
            counts['vector_shapes']+=1
        elif kind=='circle':
            r=float(a['r']);x=float(a['cx'])-r;y=float(a['cy'])-r
            shape=slide.shapes.add_shape(MSO_SHAPE.OVAL,px(x),px(y),px(2*r),px(2*r))
            fill_style(shape,a['fill']);line_style(shape,'none')
            shape.name=f'Identity_color_{counts["vector_shapes"]+1}'
            counts['vector_shapes']+=1
        elif kind=='line':
            shape=slide.shapes.add_connector(MSO_CONNECTOR.STRAIGHT,px(a['x1']),px(a['y1']),px(a['x2']),px(a['y2']))
            line_style(shape,a.get('stroke'),a.get('stroke-width',1),a.get('stroke-dasharray'))
            if 'marker-end' in a:
                end=OxmlElement('a:tailEnd');end.set('type','triangle');end.set('w','sm');end.set('len','sm')
                shape._element.spPr.get_or_add_ln().append(end)
            shape.name=f'Line_or_arrow_{counts["connectors"]+1}'
            counts['connectors']+=1
        elif kind=='text':
            text=e.text or ''
            x,y=float(a['x']),float(a['y'])
            size=float(a['font-size'])
            bold=int(a.get('font-weight','400'))>=600
            font=ImageFont.truetype(FONT_BOLD if bold else FONT,round(size*4))
            text_width=font.getlength(text)/4
            box_width=text_width+max(8,size*.6)
            anchor=a.get('text-anchor','start')
            left=x if anchor=='start' else x-box_width if anchor=='end' else x-box_width/2
            shape=slide.shapes.add_textbox(px(left),px(y-size*.94),px(box_width),px(size*1.5))
            tf=shape.text_frame
            tf.clear();tf.margin_left=tf.margin_right=tf.margin_top=tf.margin_bottom=0
            tf.word_wrap=False;tf.vertical_anchor=MSO_ANCHOR.TOP
            p=tf.paragraphs[0]
            p.alignment={'start':PP_ALIGN.LEFT,'middle':PP_ALIGN.CENTER,'end':PP_ALIGN.RIGHT}[anchor]
            p.space_before=Pt(0);p.space_after=Pt(0)
            run=p.add_run();run.text=text
            run.font.name='Arial';run.font.size=pt(size);run.font.bold=bold;run.font.color.rgb=rgb(a['fill'])
            shape.name=f'Text_{counts["text_boxes"]+1}_{text[:48]}'
            labels.append(text);counts['text_boxes']+=1
        elif kind=='image':
            encoded=a['{http://www.w3.org/1999/xlink}href'].split(',',1)[1]
            payload=base64.b64decode(encoded)
            im=Image.open(io.BytesIO(payload))
            x,y,w,h=map(float,[a['x'],a['y'],a['width'],a['height']])
            scale=min(w/im.width,h/im.height)
            iw,ih=im.width*scale,im.height*scale
            shape=slide.shapes.add_picture(io.BytesIO(payload),px(x+(w-iw)/2),px(y+(h-ih)/2),px(iw),px(ih))
            shape.name=f'Picture_{counts["pictures"]+1}_{last_asset or "replaceable_image"}'
            counts['pictures']+=1
        else:
            raise ValueError(f'Unsupported SVG object: {kind}')
    notes=[
        '使用方法：开始 → 选择 → 选择窗格，可逐项选择文字、箭头、图例和图片。',
        '更换素材：选中图片 → 图片格式 → 更改图片；按住 Shift 缩放以保持比例。',
        '正文、箭头与底色为 PPT 原生对象；人体和点云是可独立替换的透明 PNG，不能直接编辑人体表面。',
        '导出：文件 → 导出 → PDF，或另存为 PNG。生成式图不自动成为矢量图。',
        '本页为概念示意占位模板，不是实验结果；不填写真实样例指标或声称实际输出。' if concept else
        '本页使用真实预测，人物 ID 未重标注。f042 和 f050 是抽样帧，不是边界相邻帧；完整片段 IDF1=0.951。',
        '\n英文文字与中文翻译：',
    ]
    for text in labels:
        translation=layout.TRANSLATIONS.get(text)
        if translation:notes.append(f'{text}：{translation}')
        elif text.startswith('ID '):notes.append(f'{text}：预测人物身份 {text[3:]}')
        elif text.startswith('Person '):notes.append(f'{text}：人物 {text[7:]}')
        elif 's · f' in text:notes.append(f'{text}：相对于片段起点的秒数与零起始帧号')
    for shape in slide.shapes:
        # Avoid inherited Office-theme shadows on otherwise flat SVG shapes.
        sppr=shape._element.find('{http://schemas.openxmlformats.org/presentationml/2006/main}spPr')
        if sppr is not None:
            sppr.append(OxmlElement('a:effectLst'))
        for ref in shape._element.iter('{http://schemas.openxmlformats.org/drawingml/2006/main}effectRef'):
            ref.set('idx','0')
    slide.notes_slide.notes_text_frame.text='\n'.join(notes)
    prs.core_properties.title='Shot3R 可编辑 teaser：概念模板' if concept else 'Shot3R 可编辑 teaser：真实结果'
    prs.core_properties.subject='Editable text, vector annotations, and independently replaceable raster assets'
    prs.core_properties.author='Shot3R'
    prs.save(destination)
    check=Presentation(destination)
    assert len(check.slides)==1
    assert sum(s.has_text_frame and bool(s.text) for s in check.slides[0].shapes)==counts['text_boxes']
    with zipfile.ZipFile(destination) as z:
        assert z.testzip() is None
        assert 'ppt/notesSlides/notesSlide1.xml' in z.namelist()
    return {'file':destination.name,'source':str(source.relative_to(ROOT)),**counts,
            'raster_people_not_vector_meshes':not concept,'conceptual_not_result':concept}


def main():
    EDIT.mkdir(exist_ok=True)
    concept=conceptual_template()
    specs=[(ROOT/'figures/Shot3R_teaser_v1_panorama.svg','Shot3R_teaser_editable.pptx',False),
           (ROOT/'figures/Shot3R_teaser_v1_temporal.svg','Shot3R_teaser_temporal_editable.pptx',False),
           (concept,'Shot3R_concept_template.pptx',True)]
    report=[build(src,EDIT/name,is_concept) for src,name,is_concept in specs]
    (EDIT/'pptx_validation.json').write_text(json.dumps(report,indent=2,ensure_ascii=False)+'\n')
    print(json.dumps(report,indent=2,ensure_ascii=False),flush=True)


if __name__=='__main__':main()
